import os
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

def create_presentation_from_images():
    # Ask the user for the folder containing images
    folder_path = input("Enter the folder path containing images: ")

    # Check if the folder exists
    if not os.path.exists(folder_path):
        print("The folder does not exist. Please try again.")
        return

    # Create a new PowerPoint presentation
    presentation = Presentation()

    # Loop through all files in the folder
    for file_name in os.listdir(folder_path):
        # Check if the file is an image
        if file_name.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.gif', '.tiff')):
            image_path = os.path.join(folder_path, file_name)

            # Open the image to get its dimensions
            with Image.open(image_path) as img:
                img_width, img_height = img.size

            # Create a blank slide
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # Layout 6 is a blank slide

            # Get slide dimensions
            slide_width = presentation.slide_width
            slide_height = presentation.slide_height

            # Calculate the scaling factor to fill the slide
            width_ratio = slide_width / img_width
            height_ratio = slide_height / img_height
            scale_factor = max(width_ratio, height_ratio)  # Use max to fill the slide

            # Calculate the new dimensions of the image
            new_width = int(img_width * scale_factor)
            new_height = int(img_height * scale_factor)

            # Center the image on the slide
            left = (slide_width - new_width) // 2
            top = (slide_height - new_height) // 2

            # Add the image to the slide
            slide.shapes.add_picture(image_path, left, top, width=new_width, height=new_height)

    # Save the presentation
    output_path = os.path.join(folder_path, "output_presentation.pptx")
    presentation.save(output_path)
    print(f"Presentation saved as {output_path}")

# Run the script
if __name__ == "__main__":
    create_presentation_from_images()